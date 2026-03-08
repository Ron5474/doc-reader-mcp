from mcp.server.fastmcp import FastMCP
from mcp.types import TextContent, ImageContent
from pptx import Presentation
from pptx.util import Inches
import fitz # pymupdf 
import base64


mcp = FastMCP("slides-reader")

@mcp.tool()
def read_pdf(path: str) -> list:
    """Extract text from pdf file, page by page."""
    doc = fitz.open(path)
    output = []
    for i, page in enumerate(doc, start=1):
        text = page.get_text().strip()

        tables = page.find_tables()
        if tables.tables:
            for table in tables.tables:
                df = table.to_pandas()
                markdown_table = df.to_markdown(index=False)
                # Replace the raw flattened text with markdown table
                for col in df.columns:
                    for val in df[col].astype(str):
                        text = text.replace(val, "", 1)

                text = text.strip() + f"\n\n{markdown_table}"

        if text:
            output.append(TextContent(type="text", text=f"--- Page {i} ---\n{text}"))

        for img in page.get_images():
            xref = img[0]
            base_img = doc.extract_image(xref)
            img_bytes = base_img["image"]
            ext = base_img["ext"]
            b64 = base64.standard_b64encode(img_bytes).decode("utf-8")
            output.append(ImageContent(type="image", data=b64, mimeType=f"image/{ext}"))

    doc.close()
    return output

@mcp.tool()
def read_pptx(path: str) -> list:
    """Extract text from a PowerPoint file, slide by slide."""
    prs = Presentation(path)
    output = []

    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slides.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            output.append(TextContent(type="text", text=f"--- Slide {i} ---\n" + "\n".join(texts)))

        for shape in slides.shape:
            if shape.shape_type == 13:
                img = shape.image
                b64 = base64.standard_b64encode(img.blob).decode("utf-8")
                mime = img.content_type
                output.append(ImageContent(type="image", data=b64, mimeType=mime))

    return output

if __name__ == "__main__":
    mcp.run()
