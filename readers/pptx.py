from mcp.types import TextContent, ImageContent
from pptx import Presentation
import base64
import os


def read_pptx(path: str) -> list:
    """Extract text, images, and speaker notes from a PowerPoint file, slide by slide."""
    if not os.path.exists(path):
        raise ValueError(f"File not found: {path}")

    try:
        prs = Presentation(path)
    except Exception as e:
        raise ValueError(f"Could not open PPTX: {e}") from e

    output = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)

        notes = ""
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                notes = f"\n\n--- Notes ---\n{notes_text}"

        if texts or notes:
            content = f"--- Slide {i} ---\n" + "\n".join(texts) + notes
            output.append(TextContent(type="text", text=content))

        for shape in slide.shapes:
            if shape.shape_type == 13:
                img = shape.image
                b64 = base64.standard_b64encode(img.blob).decode("utf-8")
                mime = img.content_type
                output.append(ImageContent(type="image", data=b64, mimeType=mime))

    return output
