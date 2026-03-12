import pytest
from docx import Document
from pptx import Presentation
import openpyxl


@pytest.fixture
def simple_docx(tmp_path):
    doc = Document()
    doc.add_paragraph("Hello from DOCX")
    doc.add_paragraph("Second paragraph")
    path = tmp_path / "test.docx"
    doc.save(str(path))
    return str(path)


@pytest.fixture
def docx_with_table(tmp_path):
    doc = Document()
    doc.add_paragraph("Before table")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "2"
    path = tmp_path / "table.docx"
    doc.save(str(path))
    return str(path)


@pytest.fixture
def simple_xlsx(tmp_path):
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws.append(["Name", "Score"])
    ws.append(["Alice", 90])
    ws.append(["Bob", 85])
    path = tmp_path / "test.xlsx"
    wb.save(str(path))
    return str(path)


@pytest.fixture
def xlsx_multi_sheet(tmp_path):
    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "Alpha"
    ws1.append(["X", "Y"])
    ws1.append([1, 2])
    ws2 = wb.create_sheet("Beta")
    ws2.append(["P", "Q"])
    ws2.append([3, 4])
    path = tmp_path / "multi.xlsx"
    wb.save(str(path))
    return str(path)


@pytest.fixture
def pptx_with_notes(tmp_path):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Slide Title"
    slide.placeholders[1].text = "Slide body text"
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "These are speaker notes"
    path = tmp_path / "notes.pptx"
    prs.save(str(path))
    return str(path)
